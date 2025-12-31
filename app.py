import re
import io
import math
import json
import random
from collections import Counter, defaultdict
from datetime import datetime

import streamlit as st
import streamlit.components.v1 as components

# -----------------------------
# Optional deps (safe flags)
# -----------------------------
PYPDF_OK = True
try:
    from pypdf import PdfReader
except Exception:
    PYPDF_OK = False

PPTX_OK = True
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except Exception:
    PPTX_OK = False


# =========================
# Stopwords ES (compact)
# =========================
STOPWORDS_ES = {
    "a","acá","ahí","al","algo","algunos","ante","antes","apenas","aquí","así","aun","aunque",
    "bajo","bien","cada","casi","como","con","contra","cual","cuales","cuando","cuanto","de",
    "del","desde","donde","dos","el","ella","ellas","ello","ellos","en","entre","era","es",
    "esa","esas","ese","eso","esos","esta","está","están","estas","este","esto","estos","fue",
    "ha","hace","hacia","han","hasta","hay","incluso","la","las","le","les","lo","los","más",
    "me","menos","mi","mis","mucho","muy","nada","ni","no","nos","nuestra","nuestro","o","os",
    "otra","otro","para","pero","poco","por","porque","que","qué","quien","quienes","se","sea",
    "ser","si","sí","sin","sobre","solo","su","sus","también","tan","tanto","te","tener","tiene",
    "todo","todos","tu","tus","un","una","uno","unos","y","ya","yo"
}

# =========================
# Text utilities
# =========================
def clean_text(t: str) -> str:
    t = (t or "").replace("\x00", " ")
    t = re.sub(r"[ \t]+", " ", t)
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()

def normalize_token(w: str) -> str:
    w = (w or "").lower()
    w = re.sub(r"^[^\wáéíóúüñ]+|[^\wáéíóúüñ]+$", "", w, flags=re.UNICODE)
    return w

def tokenize_words(text: str):
    words = re.findall(r"[A-Za-zÁÉÍÓÚÜÑáéíóúüñ0-9]+", text or "", flags=re.UNICODE)
    out = []
    for w in words:
        w = normalize_token(w)
        if not w or w in STOPWORDS_ES or len(w) <= 2:
            continue
        out.append(w)
    return out

def split_sentences(text: str):
    text = re.sub(r"\s+", " ", (text or "")).strip()
    if not text:
        return []
    parts = re.split(r"(?<=[\.\?\!])\s+(?=[A-ZÁÉÍÓÚÜÑ0-9])", text)
    if len(parts) <= 1:
        parts = re.split(r"\s{2,}|;\s+", text)
    return [p.strip() for p in parts if p and p.strip()]

def pdf_to_text(file_bytes: bytes) -> str:
    if not PYPDF_OK:
        raise RuntimeError("Falta pypdf. Añádelo a requirements.txt.")
    bio = io.BytesIO(file_bytes)
    reader = PdfReader(bio)
    pages = []
    for p in reader.pages:
        try:
            pages.append(p.extract_text() or "")
        except Exception:
            pages.append("")
    return clean_text("\n\n".join(pages))

# =========================
# TF-IDF summary helpers
# =========================
def tfidf_sentence_scores(sentences):
    if not sentences:
        return [], {}

    sent_tokens = [tokenize_words(s) for s in sentences]
    df = Counter()
    for toks in sent_tokens:
        for t in set(toks):
            df[t] += 1

    N = max(1, len(sentences))
    idf = {t: (math.log((N + 1) / (c + 1)) + 1.0) for t, c in df.items()}

    scores = []
    for toks in sent_tokens:
        if not toks:
            scores.append(0.0)
            continue
        tf = Counter(toks)
        score = 0.0
        for t, f in tf.items():
            score += (f / len(toks)) * idf.get(t, 0.0)
        scores.append(score)

    return scores, idf

def pick_top_sentences(sentences, n):
    if not sentences:
        return []
    scores, _ = tfidf_sentence_scores(sentences)
    idx = list(range(len(sentences)))
    idx.sort(key=lambda i: scores[i], reverse=True)
    chosen = sorted(idx[:n])
    return [sentences[i] for i in chosen]

def top_keywords(text, k=16):
    toks = tokenize_words(text)
    if not toks:
        return []
    c = Counter(toks)
    ranked = sorted(c.items(), key=lambda x: (x[1], len(x[0]) >= 4, len(x[0])), reverse=True)
    out = []
    for w, _ in ranked:
        if w not in out:
            out.append(w)
        if len(out) >= k:
            break
    return out

def find_def_sentence(term, sentences):
    patt = re.compile(rf"\b{re.escape(term)}\b", re.IGNORECASE)
    candidates = [s for s in sentences if patt.search(s)]
    if not candidates:
        return None
    cues = (" es ", " son ", "se define", "consiste", "significa", "se entiende")
    def_like = [s for s in candidates if any(c in (" " + s.lower() + " ") for c in cues)]
    return (def_like[0] if def_like else candidates[0]).strip()

# =========================
# Bullet cleanup (KEY FIX)
# =========================
def normalize_bullet_text(s: str) -> str:
    s = clean_text(s)
    # sustituye bullets incrustados por separadores claros
    s = s.replace("•", " | ")
    s = re.sub(r"\s*\|\s*", " | ", s)
    s = re.sub(r"\s+", " ", s).strip()
    return s

def split_into_small_bullets(text: str, max_len=110):
    """
    Convierte frases largas (y las que llevan ' | ') en bullets cortos y legibles.
    """
    text = normalize_bullet_text(text)
    if not text:
        return []

    # 1) si hay separador interno, lo dividimos
    parts = [p.strip() for p in text.split("|") if p.strip()]

    # 2) también corta por puntuación para mejorar legibilidad
    refined = []
    for p in parts:
        if len(p) <= max_len:
            refined.append(p)
        else:
            # cortar por ". " o "; "
            chunks = re.split(r"(?<=\.)\s+|;\s+", p)
            for c in chunks:
                c = c.strip()
                if not c:
                    continue
                # si aún es larguísimo, cortar en trozos
                while len(c) > max_len:
                    cut = c.rfind(" ", 0, max_len)
                    if cut <= 40:
                        cut = max_len
                    refined.append(c[:cut].strip())
                    c = c[cut:].strip()
                if c:
                    refined.append(c)

    # limpiar duplicados cercanos
    out = []
    seen = set()
    for r in refined:
        key = r.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(r)
    return out

# =========================
# Notes generator (clean Markdown)
# =========================
def generate_notes_sections(text, detail):
    """
    Devuelve (title, sections) para usar en PPTX y en la UI.
    sections: list[ {title: str, bullets: list[str]} ]
    """
    text = clean_text(text)
    if not text:
        return "Apuntes", [{"title":"Contenido", "bullets":["(Sin contenido)"]}]

    sents = split_sentences(text)
    kws = top_keywords(text, k=16)

    if detail == "breve":
        n = 8
    elif detail == "exhaustivo":
        n = 18
    else:
        n = 12

    key_sents = pick_top_sentences(sents, n)
    glos = []
    for term in kws[:10]:
        ds = find_def_sentence(term, sents)
        if ds:
            glos.append(f"{term}: {ds}")

    title = " ".join([k.capitalize() for k in kws[:3]]) if kws else "Apuntes"

    sections = []

    # Ideas clave (bullets cortos)
    bullets = []
    for s in key_sents:
        bullets += split_into_small_bullets(s, max_len=110)
    sections.append({"title": "Ideas clave", "bullets": bullets[:24] if bullets else ["(Sin ideas detectadas)"]})

    if glos:
        gbul = []
        for g in glos:
            gbul += split_into_small_bullets(g, max_len=120)
        sections.append({"title": "Glosario", "bullets": gbul[:22]})

    # Resumen corto
    sum_sents = pick_top_sentences(sents, 10)
    sbul = []
    for s in sum_sents[:10]:
        sbul += split_into_small_bullets(s, max_len=110)
    sections.append({"title": "Resumen", "bullets": sbul[:18]})

    # Preguntas de repaso
    q = [f"¿Cómo explicarías {k} con tus palabras?" for k in (kws[:6] if kws else ["el tema"])]
    sections.append({"title": "Preguntas de repaso", "bullets": q})

    return title, sections

def sections_to_markdown(title, sections):
    md = [f"# {title}", ""]
    for sec in sections:
        md.append(f"## {sec['title']}")
        for b in sec.get("bullets", []):
            md.append(f"- {b}")
        md.append("")
    return "\n".join(md).strip()

# =========================
# Flashcards (better quality + UI)
# =========================
def generate_flashcards(text, n_cards):
    text = clean_text(text)
    sents = split_sentences(text)
    kws = top_keywords(text, k=max(22, n_cards))

    cards = []

    # definiciones
    for term in kws:
        if len(cards) >= n_cards:
            break
        ds = find_def_sentence(term, sents)
        if ds:
            back = " ".join(split_into_small_bullets(ds, max_len=120)[:2])
            cards.append({"front": f"¿Qué es {term}?", "back": back, "tag": term})

    # cloze con frases cortas
    if len(cards) < n_cards:
        pool = [s for s in sents if 8 <= len(tokenize_words(s)) <= 26]
        random.shuffle(pool)
        for s in pool:
            if len(cards) >= n_cards:
                break
            toks = tokenize_words(s)
            present = [t for t in toks if t in kws]
            if not present:
                continue
            target = present[0]
            cloze = re.sub(rf"\b{re.escape(target)}\b", "_____", s, flags=re.IGNORECASE)
            cards.append({
                "front": f"Completa: {cloze}",
                "back": f"La palabra era: {target}",
                "tag": target
            })

    return cards[:n_cards]

# =========================
# Mindmap (real render via Mermaid.js CDN)
# =========================
def generate_mindmap_mermaid(text, max_nodes=60):
    text = clean_text(text)
    if not text:
        return "mindmap\n  root((Sin contenido))"

    sents = split_sentences(text)
    kws = top_keywords(text, k=18)
    root = kws[0].capitalize() if kws else "Tema"

    co = defaultdict(Counter)
    for s in sents:
        toks = set(tokenize_words(s))
        present = [k for k in kws if k in toks]
        for a in present:
            for b in present:
                if a != b:
                    co[a][b] += 1

    lines = ["mindmap", f"  root(({root}))"]
    count = 1
    branches = kws[1:7] if len(kws) > 1 else []
    for k in branches:
        if count >= max_nodes:
            break
        lines.append(f"    {k}")
        count += 1
        for sub, _ in co[k].most_common(3):
            if count >= max_nodes:
                break
            lines.append(f"      {sub}")
            count += 1
    return "\n".join(lines)

def render_mermaid(mermaid_code: str, height: int = 620):
    """
    Render real Mermaid diagrams in Streamlit without installing anything.
    Uses Mermaid.js via CDN.
    """
    mermaid_code = (mermaid_code or "").strip()
    if not mermaid_code:
        st.info("No hay mindmap para renderizar.")
        return

    html = f"""
<!DOCTYPE html>
<html>
<head>
  <meta charset="UTF-8" />
  <script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
  <style>
    body {{ margin: 0; padding: 0; background: transparent; }}
    .wrap {{
      padding: 10px 8px 0 8px;
      font-family: system-ui, -apple-system, Segoe UI, Roboto, Helvetica, Arial;
    }}
    .card {{
      background: #ffffff;
      border: 1px solid rgba(0,0,0,0.06);
      border-radius: 14px;
      padding: 14px;
      box-shadow: 0 8px 28px rgba(0,0,0,0.06);
    }}
  </style>
</head>
<body>
  <div class="wrap">
    <div class="card">
      <div class="mermaid">
{mermaid_code}
      </div>
    </div>
  </div>
  <script>
    mermaid.initialize({{
      startOnLoad: true,
      securityLevel: "loose",
      theme: "base",
      themeVariables: {{
        fontFamily: "Inter, Segoe UI, Arial",
        primaryColor: "#ff4b4b",
        primaryTextColor: "#111111",
        lineColor: "#c9c9c9",
        tertiaryColor: "#f7f7f7"
      }}
    }});
  </script>
</body>
</html>
"""
    components.html(html, height=height, scrolling=True)

# =========================
# PPTX PRO (real improvement)
# =========================
def _rgb(hexstr):
    hexstr = hexstr.lstrip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

def _set_bg(slide, color="#ffffff"):
    # Add a full-slide rectangle background
    shp = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color)
    shp.line.fill.background()
    # send to back
    slide.shapes._spTree.remove(shp._element)
    slide.shapes._spTree.insert(2, shp._element)

def _top_bar(slide, accent="#ff4b4b"):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()

def _left_accent(slide, accent="#ff4b4b"):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()

def _title_textbox(slide, text, y=0.65, size=34, bold=True):
    tx = slide.shapes.add_textbox(Inches(0.9), Inches(y), Inches(12.2), Inches(0.9))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = _rgb("#111111")
    r.font.name = "Calibri"
    return tx

def _subtitle(slide, text, y=1.55, size=16):
    tx = slide.shapes.add_textbox(Inches(0.92), Inches(y), Inches(12.0), Inches(0.6))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = _rgb("#555555")
    r.font.name = "Calibri"
    return tx

def _footer(slide, left_text, right_text, slide_no=None):
    y = Inches(7.08)
    txl = slide.shapes.add_textbox(Inches(0.9), y, Inches(8), Inches(0.35))
    tfl = txl.text_frame
    tfl.clear()
    p = tfl.paragraphs[0]
    p.text = left_text
    p.font.size = Pt(10)
    p.font.color.rgb = _rgb("#777777")
    p.font.name = "Calibri"

    txr = slide.shapes.add_textbox(Inches(9.5), y, Inches(3.8), Inches(0.35))
    tfr = txr.text_frame
    tfr.clear()
    p2 = tfr.paragraphs[0]
    p2.text = right_text if slide_no is None else f"{right_text} • {slide_no}"
    p2.font.size = Pt(10)
    p2.font.color.rgb = _rgb("#777777")
    p2.font.name = "Calibri"
    p2.alignment = 2

def _card_box(slide, x, y, w, h, fill="#ffffff"):
    box = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = _rgb(fill)
    box.line.color.rgb = _rgb("#e6e6e6")
    return box

def _bullets_textbox(slide, bullets, x, y, w, h, font_size=20):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb("#222222")
        p.font.name = "Calibri"
    return tx

def build_pptx_pro(title, sections, accent="#ff4b4b", max_bullets=8):
    if not PPTX_OK:
        raise RuntimeError("Falta python-pptx. Añádelo a requirements.txt.")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    date_str = datetime.now().strftime("%Y-%m-%d")

    # Cover (blank layout 6 usually, but safe to use layout 0 then overwrite)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "#fbfbfd")
    _top_bar(slide, accent)
    _left_accent(slide, accent)
    _title_textbox(slide, title, y=1.55, size=42, bold=True)
    _subtitle(slide, "Presentación automática (sin API) • Flashcards • Mindmap", y=2.45, size=16)
    _footer(slide, "StudyWave PRO", date_str)

    # Agenda
    slide_no = 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "#ffffff")
    _top_bar(slide, accent)
    _title_textbox(slide, "Agenda", y=0.85, size=34)
    agenda = [s["title"] for s in sections][:12]
    _card_box(slide, 0.9, 1.7, 12.0, 4.9, fill="#ffffff")
    numbered = [f"{i+1}. {t}" for i, t in enumerate(agenda)]
    _bullets_textbox(slide, numbered, 1.15, 1.95, 11.5, 4.4, font_size=22)
    _footer(slide, "Agenda", date_str, slide_no)
    slide_no += 1

    # Content
    for sec in sections:
        bullets_raw = sec.get("bullets", []) or ["(Sin puntos)"]

        # recortar y limpiar
        bullets = []
        for b in bullets_raw:
            bullets += split_into_small_bullets(b, max_len=110)
        bullets = bullets[:60] if bullets else ["(Sin puntos)"]

        # Section divider
        div = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(div, "#fbfbfd")
        _top_bar(div, accent)
        _left_accent(div, accent)
        _title_textbox(div, sec["title"], y=2.3, size=46, bold=True)
        _subtitle(div, "Bloque de contenidos", y=3.25, size=16)
        _footer(div, sec["title"], date_str, slide_no)
        slide_no += 1

        # Slides for bullets
        chunks = [bullets[i:i+max_bullets*2] for i in range(0, len(bullets), max_bullets*2)]
        for chunk in chunks:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_bg(slide, "#ffffff")
            _top_bar(slide, accent)
            _title_textbox(slide, sec["title"], y=0.7, size=32)

            # Two columns if many bullets
            if len(chunk) > max_bullets:
                left = chunk[:max_bullets]
                right = chunk[max_bullets:max_bullets*2]

                _card_box(slide, 0.9, 1.55, 5.95, 5.25, fill="#ffffff")
                _card_box(slide, 6.95, 1.55, 5.95, 5.25, fill="#ffffff")
                _bullets_textbox(slide, left, 1.15, 1.8, 5.45, 4.8, font_size=20)
                _bullets_textbox(slide, right, 7.2, 1.8, 5.45, 4.8, font_size=20)
            else:
                _card_box(slide, 0.9, 1.55, 12.0, 5.25, fill="#ffffff")
                _bullets_textbox(slide, chunk, 1.15, 1.8, 11.5, 4.8, font_size=22)

            _footer(slide, sec["title"], date_str, slide_no)
            slide_no += 1

    # Final takeaways
    all_bul = []
    for s in sections:
        for b in s.get("bullets", []):
            all_bul += split_into_small_bullets(b, max_len=110)
    take = all_bul[:10] if all_bul else ["(Sin conclusiones)"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, "#ffffff")
    _top_bar(slide, accent)
    _title_textbox(slide, "Takeaways", y=0.8, size=34)
    _card_box(slide, 0.9, 1.55, 12.0, 5.25, fill="#ffffff")
    _bullets_textbox(slide, take, 1.15, 1.8, 11.5, 4.8, font_size=22)
    _footer(slide, "Cierre", date_str, slide_no)

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()

# =========================
# Flashcards UI (real cards)
# =========================
def inject_card_css():
    st.markdown(
        """
<style>
.sw-grid { display: grid; grid-template-columns: 1fr; gap: 14px; }
@media (min-width: 900px) { .sw-grid { grid-template-columns: 1fr 1fr; } }
.sw-card {
  background: #ffffff;
  border: 1px solid rgba(0,0,0,0.06);
  border-radius: 18px;
  padding: 16px 16px;
  box-shadow: 0 10px 28px rgba(0,0,0,0.06);
}
.sw-pill {
  display: inline-block;
  font-size: 12px;
  padding: 4px 10px;
  border-radius: 999px;
  background: rgba(255,75,75,0.12);
  color: #b61f1f;
  border: 1px solid rgba(255,75,75,0.25);
  margin-bottom: 10px;
}
.sw-front { font-size: 18px; font-weight: 700; color: #111; margin-bottom: 10px; }
.sw-back  { font-size: 16px; color: #222; line-height: 1.35; }
.sw-muted { color: #666; font-size: 12px; }
</style>
""",
        unsafe_allow_html=True
    )

def flashcard_view(cards):
    if not cards:
        st.info("No hay flashcards generadas.")
        return

    if "fc_idx" not in st.session_state:
        st.session_state.fc_idx = 0
    if "fc_show_back" not in st.session_state:
        st.session_state.fc_show_back = False

    idx = st.session_state.fc_idx % len(cards)
    card = cards[idx]

    inject_card_css()

    st.markdown('<div class="sw-card">', unsafe_allow_html=True)
    st.markdown(f'<div class="sw-pill">#{card.get("tag","")}</div>', unsafe_allow_html=True)
    st.markdown(f'<div class="sw-front">{card.get("front","")}</div>', unsafe_allow_html=True)

    if st.session_state.fc_show_back:
        st.markdown(f'<div class="sw-back">{card.get("back","")}</div>', unsafe_allow_html=True)
    else:
        st.markdown('<div class="sw-muted">Pulsa “Mostrar respuesta” para ver el reverso.</div>', unsafe_allow_html=True)

    st.markdown(f'<div class="sw-muted">Tarjeta {idx+1} / {len(cards)}</div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

    c1, c2, c3, c4 = st.columns([1,1,1,1])
    with c1:
        if st.button("⬅️ Anterior", use_container_width=True):
            st.session_state.fc_idx = (idx - 1) % len(cards)
            st.session_state.fc_show_back = False
            st.rerun()
    with c2:
        if st.button("🔁 Barajar", use_container_width=True):
            random.shuffle(cards)
            st.session_state.fc_idx = 0
            st.session_state.fc_show_back = False
            st.rerun()
    with c3:
        if st.button("👁️ Mostrar respuesta" if not st.session_state.fc_show_back else "🙈 Ocultar respuesta", use_container_width=True):
            st.session_state.fc_show_back = not st.session_state.fc_show_back
            st.rerun()
    with c4:
        if st.button("Siguiente ➡️", use_container_width=True):
            st.session_state.fc_idx = (idx + 1) % len(cards)
            st.session_state.fc_show_back = False
            st.rerun()

    st.progress((idx + 1) / len(cards))

def flashcards_export_json(cards):
    return {"flashcards": [{"front": c["front"], "back": c["back"], "tags":[c.get("tag","")]} for c in cards]}

def flashcards_to_anki_tsv(cards):
    lines = []
    for c in cards:
        front = (c.get("front","") or "").replace("\t"," ").strip()
        back = (c.get("back","") or "").replace("\t"," ").strip()
        tag = c.get("tag","") or ""
        lines.append(f"{front}\t{back}\t{tag}")
    return "\n".join(lines)

# =========================
# Streamlit UI
# =========================
st.set_page_config(page_title="StudyWave PRO (sin API)", page_icon="🧠", layout="wide")
st.title("🧠 StudyWave PRO — PPTX + Flashcards + Mindmap (SIN API, sin modelos)")

with st.sidebar:
    st.subheader("Estado librerías")
    st.write(f"pypdf (PDF): {'✅' if PYPDF_OK else '❌'}")
    st.write(f"python-pptx (PPTX): {'✅' if PPTX_OK else '❌'}")

    st.divider()
    source = st.radio("Fuente", ["PDF", "Texto"], index=0)
    detail = st.selectbox("Nivel de apuntes", ["breve", "medio", "exhaustivo"], index=1)

    st.divider()
    accent = st.color_picker("Color acento", "#ff4b4b")
    bullets_per_col = st.slider("Bullets por columna (PPT)", 5, 10, 8, 1)
    n_flash = st.slider("Nº flashcards", 8, 60, 24, 1)

    st.divider()
    st.caption("Sin API. No se instala IA. Todo se hace con reglas + NLP clásica.")

content = ""
doc_name = "apuntes"

if source == "Texto":
    content = st.text_area("Pega tus apuntes", height=280, placeholder="Pega aquí tus apuntes o texto...")
else:
    pdf = st.file_uploader("Sube un PDF", type=["pdf"])
    if pdf is not None:
        doc_name = re.sub(r"\.pdf$", "", pdf.name, flags=re.IGNORECASE) or "apuntes"
        try:
            if not PYPDF_OK:
                st.error("No puedo leer PDF porque falta pypdf en requirements.txt.")
                content = ""
            else:
                content = pdf_to_text(pdf.read())
                st.success(f"Texto extraído: {len(content):,} caracteres")
                if len(content.strip()) < 300:
                    st.warning("He extraído muy poco texto. Si tu PDF es escaneado (imagen), sin OCR no hay texto real.")
        except Exception as e:
            st.error("Error leyendo el PDF:")
            st.exception(e)
            content = ""

st.divider()
go = st.button("✨ Generar materiales PRO", type="primary", use_container_width=True)

if go:
    try:
        if not clean_text(content):
            st.error("No hay texto para procesar.")
            st.stop()

        title, sections = generate_notes_sections(content, detail)
        notes_md = sections_to_markdown(title, sections)

        # Generate components
        cards = generate_flashcards(content, n_flash)
        mindmap = generate_mindmap_mermaid(content, max_nodes=70)

        # Layout tabs
        tab1, tab2, tab3, tab4 = st.tabs(["📚 Apuntes", "🃏 Flashcards (visual)", "🧩 Mindmap (real)", "📊 PowerPoint PRO"])

        with tab1:
            st.subheader("Apuntes (limpios y listos)")
            st.markdown(notes_md)
            st.download_button(
                "⬇️ Descargar apuntes (.md)",
                data=notes_md.encode("utf-8"),
                file_name=f"{doc_name}_apuntes.md",
                mime="text/markdown",
                use_container_width=True,
            )

        with tab2:
            st.subheader("Flashcards (tipo tarjetas)")
            flashcard_view(cards)

            st.divider()
            c1, c2 = st.columns(2)
            with c1:
                fc_json = flashcards_export_json(cards)
                st.download_button(
                    "⬇️ Descargar flashcards (.json)",
                    data=json.dumps(fc_json, ensure_ascii=False, indent=2).encode("utf-8"),
                    file_name=f"{doc_name}_flashcards.json",
                    mime="application/json",
                    use_container_width=True,
                )
            with c2:
                anki = flashcards_to_anki_tsv(cards)
                st.download_button(
                    "⬇️ Descargar para Anki (.tsv)",
                    data=anki.encode("utf-8"),
                    file_name=f"{doc_name}_anki.tsv",
                    mime="text/tab-separated-values",
                    use_container_width=True,
                )

        with tab3:
            st.subheader("Mapa mental renderizado (no código)")
            render_mermaid(mindmap, height=680)

            st.download_button(
                "⬇️ Descargar mindmap (Mermaid .mmd)",
                data=mindmap.encode("utf-8"),
                file_name=f"{doc_name}_mindmap.mmd",
                mime="text/plain",
                use_container_width=True,
            )

        with tab4:
            st.subheader("PowerPoint PRO")
            if not PPTX_OK:
                st.error("No puedo generar PPTX porque falta python-pptx en requirements.txt.")
            else:
                pptx_bytes = build_pptx_pro(
                    title=doc_name,
                    sections=sections,
                    accent=accent,
                    max_bullets=bullets_per_col
                )
                st.success(f"PPTX generado ✅ ({len(pptx_bytes)/1024:.1f} KB)")
                st.download_button(
                    "⬇️ Descargar PowerPoint (.pptx)",
                    data=pptx_bytes,
                    file_name=f"{doc_name}_presentacion_PRO.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                )

        st.success("Listo ✅ (PPTX + Flashcards visuales + Mindmap real)")

    except Exception as e:
        st.error("Ocurrió un error (capturado para que la app no se caiga). Detalle:")
        st.exception(e)
